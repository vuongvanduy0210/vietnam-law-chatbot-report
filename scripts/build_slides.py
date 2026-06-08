# -*- coding: utf-8 -*-
"""
Slide báo cáo bảo vệ: Trợ lý ảo Tư vấn Pháp luật Việt Nam (Agentic RAG).
Phong cách hiện đại, tone Indigo + đa sắc nhấn, nền gradient sáng.
Hiệu ứng: transition fade + entrance float-up/fade gợn sóng (tự chạy khi vào slide) + wipe cho connector.
Chạy: /tmp/pptxenv/bin/python build_slides.py
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn
from lxml import etree

# ----------------------------------------------------------------------------
COVER = {
    "title": "TRỢ LÝ ẢO TƯ VẤN PHÁP LUẬT VIỆT NAM",
    "subtitle": "Hệ thống hỏi – đáp pháp luật theo kiến trúc Agentic RAG",
    "student": "Vương Văn Duy",
    "advisor": "[Giảng viên hướng dẫn — đang cập nhật]",
    "school": "[Trường / Khoa — đang cập nhật]",
    "occasion": "Đồ án tốt nghiệp · 2026",
}
PROJECT_FOOTER = "Trợ lý ảo Tư vấn Pháp luật Việt Nam · Agentic RAG"
IMG_DIR = "/Users/duy/Downloads/sourcecode/Bao_Cao/sample_images/chapter4"
OUT = "/Users/duy/Downloads/sourcecode/Bao_Cao/Slide_Bao_Cao.pptx"

# ----------------------------------------------------------------------------
# MÀU (RGBColor cho text/line; hex string cho gradient)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
SLATE   = RGBColor(0x0F, 0x17, 0x2A)
SLATE2  = RGBColor(0x33, 0x41, 0x55)
MUTED   = RGBColor(0x6B, 0x72, 0x80)
BORDER  = RGBColor(0xE5, 0xE9, 0xF0)
INDIGO  = RGBColor(0x4F, 0x46, 0xE5)
INDIGO_D= RGBColor(0x37, 0x30, 0xA3)
SKY     = RGBColor(0x0E, 0xA5, 0xE9)
TEAL    = RGBColor(0x0D, 0x94, 0x88)
AMBER   = RGBColor(0xD9, 0x77, 0x06)
VIOLET  = RGBColor(0x7C, 0x3A, 0xED)
ROSE    = RGBColor(0xE1, 0x1D, 0x48)
GREEN   = RGBColor(0x0F, 0x9D, 0x58)
CLOUD   = RGBColor(0xF8, 0xFA, 0xFF)

# hex
H_WHITE="FFFFFF"; H_BG1="FFFFFF"; H_BG2="EEF1FE"; H_BG3="E9F4FE"
H_INDIGO="4F46E5"; H_INDIGO_D="3730A3"; H_SKY="0EA5E9"; H_VIOLET="7C3AED"
H_TEAL="0D9488"; H_AMBER="F59E0B"; H_GREEN="10B981"; H_ROSE="F43F5E"; H_SLATE="0F172A"

# (accent RGB, light-tint hex) cho card đa sắc
PAL = [
    (INDIGO,  "EEEDFE"),
    (SKY,     "E8F6FE"),
    (TEAL,    "E5F6F3"),
    (AMBER,   "FDF3E2"),
    (VIOLET,  "F2EBFE"),
    (ROSE,    "FDEAEF"),
]

FONT = "Calibri"
FONT_L = "Calibri Light"
EMU_IN = 914400
SW, SH = 13.333, 7.5

prs = Presentation()
prs.slide_width = Inches(SW)
prs.slide_height = Inches(SH)
BLANK = prs.slide_layouts[6]
DECOR = set()   # shape_id của hình trang trí → không animate


# ============================================================================
# LOW-LEVEL HELPERS
# ============================================================================
def set_grad(shape, hex1, hex2, angle=90):
    """Đổi fill của shape thành gradient tuyến tính hex1→hex2 (góc độ, từ trục ngang)."""
    spPr = shape._element.spPr
    idx = None
    for tag in ("a:solidFill", "a:gradFill", "a:noFill", "a:blipFill", "a:pattFill", "a:grpFill"):
        el = spPr.find(qn(tag))
        if el is not None:
            idx = list(spPr).index(el); spPr.remove(el); break
    gf = etree.Element(qn("a:gradFill"))
    gsLst = etree.SubElement(gf, qn("a:gsLst"))
    for pos, hx in ((0, hex1), (100000, hex2)):
        gs = etree.SubElement(gsLst, qn("a:gs")); gs.set("pos", str(pos))
        etree.SubElement(gs, qn("a:srgbClr")).set("val", hx)
    lin = etree.SubElement(gf, qn("a:lin")); lin.set("ang", str(int(angle * 60000))); lin.set("scaled", "1")
    if idx is not None:
        spPr.insert(idx, gf)
    else:
        ln = spPr.find(qn("a:ln"))
        spPr.insert(list(spPr).index(ln), gf) if ln is not None else spPr.append(gf)


def _soft_shadow(sp, blur=110000, dist=46000, alpha=18000):
    spPr = sp._element.spPr
    eff = spPr.find(qn("a:effectLst"))
    if eff is None:
        eff = etree.SubElement(spPr, qn("a:effectLst"))
    sh = etree.SubElement(eff, qn("a:outerShdw"))
    sh.set("blurRad", str(blur)); sh.set("dist", str(dist)); sh.set("dir", "5400000"); sh.set("rotWithShape", "0")
    c = etree.SubElement(sh, qn("a:srgbClr")); c.set("val", "1E293B")
    etree.SubElement(c, qn("a:alpha")).set("val", str(alpha))


def rrect(s, x, y, w, h, fill=WHITE, line=BORDER, line_w=1.0, radius=0.09, shadow=False,
          grad=None, grad_ang=90):
    sp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    try: sp.adjustments[0] = radius
    except Exception: pass
    if grad: set_grad(sp, grad[0], grad[1], grad_ang)
    if shadow: _soft_shadow(sp)
    return sp


def rect(s, x, y, w, h, fill, grad=None, grad_ang=0):
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.shadow.inherit = False
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    sp.line.fill.background()
    if grad: set_grad(sp, grad[0], grad[1], grad_ang)
    return sp


def oval(s, x, y, d, fill=None, grad=None, grad_ang=90, line=None, decor=False):
    sp = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d if isinstance(d, (int, float)) else d))
    sp.shadow.inherit = False
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb = line; sp.line.width = Pt(2)
    if grad: set_grad(sp, grad[0], grad[1], grad_ang)
    if decor: DECOR.add(sp.shape_id)
    return sp


def tb(s, x, y, w, h, lines, anchor=MSO_ANCHOR.TOP):
    box = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = ln.get("align", PP_ALIGN.LEFT)
        if ln.get("space_after") is not None: p.space_after = Pt(ln["space_after"])
        if ln.get("space_before") is not None: p.space_before = Pt(ln["space_before"])
        if ln.get("line_spacing"): p.line_spacing = ln["line_spacing"]
        runs = ln["text"] if isinstance(ln["text"], list) else [ln]
        for rspec in runs:
            r = p.add_run(); r.text = rspec["text"]
            r.font.size = Pt(rspec.get("size", ln.get("size", 14)))
            r.font.bold = rspec.get("bold", ln.get("bold", False))
            r.font.name = rspec.get("font", ln.get("font", FONT))
            r.font.color.rgb = rspec.get("color", ln.get("color", SLATE))
    return box


def _center_text(shape, text, size, color, bold=True, font=FONT):
    tf = shape.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.04); tf.margin_right = Inches(0.04); tf.margin_top = 0; tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text; r.font.size = Pt(size); r.font.bold = bold
    r.font.name = font; r.font.color.rgb = color


def pill(s, x, y, text, fill_hex_a, fill_hex_b, color=WHITE, h=0.42, size=12, pad=0.34):
    w = pad + len(text) * size * 0.0098
    p = rrect(s, x, y, w, h, fill=INDIGO, line=None, radius=0.5, grad=(fill_hex_a, fill_hex_b), grad_ang=0)
    _center_text(p, text, size, color)
    return p, w


def icon(s, x, y, d, emoji, a_hex, b_hex, size=18, radius_oval=True):
    sp = oval(s, x, y, d, fill=INDIGO, grad=(a_hex, b_hex), grad_ang=45)
    _soft_shadow(sp, blur=60000, dist=24000, alpha=22000)
    _center_text(sp, emoji, size, WHITE)
    return sp


def num_badge(s, idx, n=13):
    d = 0.66
    sp = oval(s, SW - 1.15, 0.42, d, fill=INDIGO, grad=(H_INDIGO, H_VIOLET), grad_ang=45)
    _soft_shadow(sp, blur=60000, dist=20000, alpha=22000)
    tf = sp.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = str(idx); r.font.size = Pt(18); r.font.bold = True; r.font.name = FONT; r.font.color.rgb = WHITE
    r2 = p.add_run(); r2.text = f"/{n}"; r2.font.size = Pt(10); r2.font.bold = True; r2.font.name = FONT
    r2.font.color.rgb = RGBColor(0xE0, 0xE7, 0xFF)
    return sp


def slide(decor_style="corner"):
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.shadow.inherit = False; bg.fill.solid(); bg.fill.fore_color.rgb = WHITE; bg.line.fill.background()
    set_grad(bg, H_BG1, H_BG2, 115)
    DECOR.add(bg.shape_id)
    if decor_style == "corner":
        o1 = oval(s, 10.4, -2.0, 4.6, grad=("EAECFE", "F3F0FE"), grad_ang=30, decor=True)
        o2 = oval(s, -1.6, 5.2, 3.4, grad=("E7F5FE", "EFFAFE"), grad_ang=30, decor=True)
    return s


def header(s, kicker, title, idx, n=13, accent=(H_INDIGO, H_VIOLET)):
    # dải gradient mảnh trên cùng
    band = rect(s, 0, 0, SW, 0.16, INDIGO, grad=(H_INDIGO, H_SKY), grad_ang=0)
    pill(s, 0.7, 0.45, kicker.upper(), accent[0], accent[1], h=0.4, size=11.5)
    tb(s, 0.7, 0.95, 11.0, 0.72, [{"text": title, "size": 27, "bold": True, "color": SLATE}])
    u = rrect(s, 0.72, 1.66, 1.5, 0.08, fill=INDIGO, line=None, radius=0.5, grad=(H_INDIGO, H_SKY), grad_ang=0)
    num_badge(s, idx, n)
    footer(s, idx, n)


def footer(s, idx, n=13):
    rect(s, 0.7, 7.0, 12.0, 0.012, BORDER)
    tb(s, 0.7, 7.06, 9.0, 0.3, [{"text": PROJECT_FOOTER, "size": 9, "color": MUTED}])
    tb(s, 11.5, 7.06, 1.1, 0.3, [{"text": f"{idx:02d}", "size": 9, "bold": True, "color": INDIGO, "align": PP_ALIGN.RIGHT}])


def bullets(s, x, y, w, h, items, size=14, gap=10, color=SLATE2, mk=INDIGO, lh=1.14):
    box = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = 0; tf.margin_top = 0; tf.margin_right = 0; tf.margin_bottom = 0
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap); p.line_spacing = lh
        rm = p.add_run(); rm.text = "●  "; rm.font.size = Pt(size * 0.62); rm.font.bold = True
        rm.font.name = FONT; rm.font.color.rgb = mk
        if isinstance(it, tuple):
            rb = p.add_run(); rb.text = it[0]; rb.font.size = Pt(size); rb.font.bold = True
            rb.font.name = FONT; rb.font.color.rgb = SLATE
            rr = p.add_run(); rr.text = it[1]; rr.font.size = Pt(size); rr.font.name = FONT; rr.font.color.rgb = color
        else:
            rr = p.add_run(); rr.text = it; rr.font.size = Pt(size); rr.font.name = FONT; rr.font.color.rgb = color
    return box


def box_node(s, x, y, w, h, title, sub=None, accent=INDIGO, a_hex=H_INDIGO, tint="FFFFFF",
             icon_emoji=None, line_w=1.6):
    sp = rrect(s, x, y, w, h, fill=WHITE, line=accent, line_w=line_w, radius=0.14, shadow=True,
               grad=("FFFFFF", tint), grad_ang=90)
    top = y + 0.14
    if icon_emoji:
        icon(s, x + 0.18, y + 0.16, 0.46, icon_emoji, a_hex, a_hex, size=15)
        tx = x + 0.74; tw = w - 0.9
    else:
        tx = x; tw = w
    lines = [{"text": title, "size": 12.5, "bold": True, "color": accent,
              "align": (PP_ALIGN.LEFT if icon_emoji else PP_ALIGN.CENTER), "space_after": 1}]
    if sub:
        lines.append({"text": sub, "size": 9.3, "color": MUTED,
                      "align": (PP_ALIGN.LEFT if icon_emoji else PP_ALIGN.CENTER)})
    bx = s.shapes.add_textbox(Inches(tx), Inches(y), Inches(tw), Inches(h))
    tf = bx.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.04); tf.margin_right = Inches(0.04)
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = ln["align"]
        if ln.get("space_after"): p.space_after = Pt(ln["space_after"])
        r = p.add_run(); r.text = ln["text"]; r.font.size = Pt(ln["size"]); r.font.bold = ln.get("bold", False)
        r.font.name = FONT; r.font.color.rgb = ln["color"]
    return sp


def connect(s, x1, y1, x2, y2, color=INDIGO, w=2.0, arrow=True, kind=MSO_CONNECTOR.STRAIGHT):
    cn = s.shapes.add_connector(kind, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    cn.line.color.rgb = color; cn.line.width = Pt(w); cn.shadow.inherit = False
    if arrow:
        ln = cn.line._get_or_add_ln()
        e = etree.SubElement(ln, qn("a:tailEnd")); e.set("type", "triangle"); e.set("w", "med"); e.set("len", "med")
    return cn


def pic_h(s, x, y, h, path, caption):
    if not os.path.exists(path): return None
    pic = s.shapes.add_picture(path, Inches(x), Inches(y), height=Inches(h))
    pic.line.color.rgb = BORDER; pic.line.width = Pt(1.0); _soft_shadow(pic, blur=70000, dist=24000, alpha=16000)
    w_in = pic.width / EMU_IN
    tb(s, x - 0.4, y + h + 0.05, w_in + 0.8, 0.3, [{"text": caption, "size": 9.5, "color": MUTED, "align": PP_ALIGN.CENTER}])
    return pic


def pic_w(s, x, y, w, path, caption):
    if not os.path.exists(path): return None
    pic = s.shapes.add_picture(path, Inches(x), Inches(y), width=Inches(w))
    pic.line.color.rgb = BORDER; pic.line.width = Pt(1.0); _soft_shadow(pic, blur=70000, dist=24000, alpha=16000)
    h_in = pic.height / EMU_IN
    tb(s, x, y + h_in + 0.04, w, 0.3, [{"text": caption, "size": 9.5, "color": MUTED, "align": PP_ALIGN.CENTER}])
    return pic


# ============================================================================
# SLIDE 1 — TRANG BÌA
# ============================================================================
s = slide(decor_style="none")
bg = s.shapes[0]  # đã là gradient
# hero panel phải
panel = rrect(s, 8.55, -0.6, 5.6, 8.7, fill=INDIGO, line=None, radius=0.0, grad=(H_INDIGO_D, H_VIOLET), grad_ang=125)
DECOR.add(panel.shape_id)
o = oval(s, 9.7, 1.0, 3.9, grad=("6D5DF0", "8B5CF6"), grad_ang=40, decor=True)
o2 = oval(s, 11.6, 4.7, 2.6, grad=("4338CA", "6366F1"), grad_ang=40, decor=True)
scale = oval(s, 10.55, 2.05, 1.9, fill=WHITE, decor=False); scale.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
_center_text(scale, "⚖️", 54, INDIGO_D)
DECOR.add(scale.shape_id)
# dải accent trái
rect(s, 0, 0, 0.26, SH, INDIGO, grad=(H_INDIGO, H_SKY), grad_ang=90)
pill(s, 0.9, 0.95, "ĐỒ ÁN TỐT NGHIỆP · 2026", H_INDIGO, H_VIOLET, h=0.48, size=12.5)
tb(s, 0.88, 1.95, 7.7, 2.3, [{"text": COVER["title"], "size": 38, "bold": True, "color": SLATE, "line_spacing": 1.05}])
rrect(s, 0.92, 4.05, 1.7, 0.09, fill=INDIGO, line=None, radius=0.5, grad=(H_INDIGO, H_SKY), grad_ang=0)
tb(s, 0.9, 4.3, 7.4, 0.9, [{"text": COVER["subtitle"], "size": 17, "color": INDIGO_D, "font": FONT_L, "line_spacing": 1.12}])
tb(s, 0.9, 5.55, 7.4, 1.5, [
    {"text": [{"text": "Sinh viên:  ", "size": 13.5, "color": MUTED}, {"text": COVER["student"], "size": 13.5, "bold": True, "color": SLATE}], "space_after": 7},
    {"text": [{"text": "GVHD:  ", "size": 13.5, "color": MUTED}, {"text": COVER["advisor"], "size": 13.5, "bold": True, "color": SLATE}], "space_after": 7},
    {"text": [{"text": "Đơn vị:  ", "size": 13.5, "color": MUTED}, {"text": COVER["school"], "size": 13.5, "bold": True, "color": SLATE}]},
])
# chips công nghệ nổi bật
for i, t in enumerate(["Agentic RAG", "LangGraph", "Gemini", "Vector Search"]):
    pl, w = pill(s, 0.9 + i * 1.0, 6.95, t, "EEF1FE", "EAF6FE", color=INDIGO_D, h=0.34, size=9.5, pad=0.22)

# ============================================================================
# SLIDE 2 — LÝ DO CHỌN ĐỀ TÀI
# ============================================================================
s = slide()
header(s, "Bối cảnh & Vấn đề", "Vì sao cần một trợ lý pháp luật đáng tin?", 2)
cards = [
    ("📚", "Nhu cầu lớn, tiếp cận khó", "Văn bản pháp luật đồ sộ, ngôn ngữ chuyên ngành; người dân khó tự tra cứu đúng điều khoản áp dụng.", INDIGO, H_INDIGO, "EEEDFE"),
    ("🤖", "LLM thuần dễ “bịa”", "Chatbot thường tự tin nêu sai số điều, mức phạt, thời hạn — rủi ro nghiêm trọng trong pháp lý.", AMBER, H_AMBER, "FDF3E2"),
    ("🔄", "Pháp luật thay đổi liên tục", "Quy định đúng hôm qua có thể đã bị thay thế. VD: Nghị định 168/2024 thay 100/2019 về giao thông.", SKY, H_SKY, "E8F6FE"),
]
cw, gap, x0, cy, ch = 3.82, 0.24, 0.7, 1.98, 2.72
for i, (em, t, d, col, hx, tint) in enumerate(cards):
    x = x0 + i * (cw + gap)
    rrect(s, x, cy, cw, ch, line=BORDER, radius=0.1, shadow=True, grad=("FFFFFF", tint), grad_ang=120)
    rrect(s, x, cy, cw, 0.13, fill=col, line=None, radius=0.0)
    icon(s, x + 0.32, cy + 0.36, 0.62, em, hx, hx, size=20)
    tb(s, x + 0.32, cy + 1.12, cw - 0.6, 0.6, [{"text": t, "size": 15, "bold": True, "color": SLATE}])
    tb(s, x + 0.32, cy + 1.62, cw - 0.6, 1.0, [{"text": d, "size": 11.5, "color": SLATE2, "line_spacing": 1.14}])
band = rrect(s, 0.7, 5.0, 11.96, 1.32, fill=INDIGO, line=None, radius=0.12, shadow=True, grad=(H_SLATE, H_INDIGO_D), grad_ang=0)
tb(s, 1.05, 5.16, 11.3, 1.05, [
    {"text": [{"text": "→ Cần hệ thống trả lời ", "size": 16, "color": WHITE},
              {"text": "CÓ CĂN CỨ", "size": 16, "bold": True, "color": RGBColor(0x7D,0xD3,0xFC)},
              {"text": " · biết ", "size": 16, "color": WHITE},
              {"text": "KIỂM CHỨNG HIỆU LỰC", "size": 16, "bold": True, "color": RGBColor(0x7D,0xD3,0xFC)},
              {"text": " · ", "size": 16, "color": WHITE},
              {"text": "CHỐNG BỊA ĐẶT", "size": 16, "bold": True, "color": RGBColor(0x7D,0xD3,0xFC)}], "space_after": 4, "line_spacing": 1.12},
    {"text": "Đây chính là bài toán mà đề tài lựa chọn giải quyết.", "size": 12, "color": RGBColor(0xCB,0xD5,0xE1)},
], anchor=MSO_ANCHOR.MIDDLE)

# ============================================================================
# SLIDE 3 — MỤC TIÊU & PHẠM VI
# ============================================================================
s = slide()
header(s, "Mục tiêu & Phạm vi", "Mục tiêu và phạm vi đề tài", 3)
rrect(s, 0.7, 1.98, 7.5, 4.45, line=BORDER, radius=0.06, shadow=True, grad=("FFFFFF", "F7F8FF"), grad_ang=120)
icon(s, 1.0, 2.22, 0.5, "🎯", H_INDIGO, H_VIOLET, size=16)
tb(s, 1.62, 2.28, 6.4, 0.5, [{"text": "MỤC TIÊU", "size": 13.5, "bold": True, "color": INDIGO_D}])
bullets(s, 1.0, 2.95, 7.0, 3.4, [
    ("Hỏi – đáp pháp luật tiếng Việt ", "trả lời dựa trên nguồn, có trích dẫn điều khoản rõ ràng."),
    ("Tự động kiểm chứng tính cập nhật ", "qua tra cứu web + chống trả lời bịa đặt bằng bước kiểm chứng."),
    ("Số hoá & quản trị văn bản: ", "admin upload PDF → tự trích xuất điều luật → nạp kho tri thức."),
    ("Đa nền tảng: ", "ứng dụng di động cho người dùng và web quản trị cho cán bộ."),
], size=13.5, gap=13)
rrect(s, 8.4, 1.98, 4.26, 4.45, radius=0.08, line=None, grad=(H_SKY, H_TEAL), grad_ang=130, shadow=True)
icon(s, 8.7, 2.24, 0.5, "🧭", "FFFFFF", "E8F6FE", size=16)
icn = s.shapes[-1]
tb(s, 9.32, 2.3, 3.1, 0.5, [{"text": "PHẠM VI", "size": 13.5, "bold": True, "color": WHITE}])
bullets(s, 8.7, 3.0, 3.7, 3.3, [
    "Lĩnh vực: giao thông, dân sự, lao động, hành chính…",
    "Dữ liệu pháp luật nội bộ cập nhật đến năm 2026.",
    "Trả lời mang tính tham khảo, có tuyên bố miễn trừ.",
], size=12.5, gap=12, color=WHITE, mk=WHITE)

# ============================================================================
# SLIDE 4 — PHƯƠNG ÁN: AGENTIC RAG
# ============================================================================
s = slide()
header(s, "Phương án giải quyết", "Agentic RAG — vì sao không để LLM tự trả lời?", 4)
cols = [
    ("LLM thuần", ["Trả lời nhanh", "Không nguồn, dễ bịa", "Không biết luật mới"], MUTED, "9CA3AF", "F3F4F6", False),
    ("RAG cơ bản", ["Có nguồn nội bộ", "Dữ liệu có thể đã cũ", "Không kiểm chứng hiệu lực"], SKY, H_SKY, "E8F6FE", False),
    ("Agentic RAG  ·  đề tài", ["Bắt buộc 2 nguồn:\nnội bộ + web hiện hành", "Agent điều phối công cụ", "Verifier kiểm chứng,\nchống bịa đặt"], INDIGO, H_INDIGO, "EEEDFE", True),
]
cw, gap, x0, cy, ch = 3.82, 0.24, 0.7, 1.98, 3.05
for i, (t, items, col, hx, tint, hi) in enumerate(cols):
    x = x0 + i * (cw + gap)
    rrect(s, x, cy, cw, ch, fill=WHITE, line=(col if hi else BORDER), line_w=(2.4 if hi else 1.0),
          radius=0.1, shadow=True, grad=("FFFFFF", tint), grad_ang=120)
    cap = rrect(s, x, cy, cw, 0.66, fill=col, line=None, radius=0.1, grad=(hx, hx), grad_ang=0)
    _center_text(cap, t, 14, WHITE)
    yy = cy + 0.92
    for it in items:
        mk = "✓  " if hi else "•  "
        tb(s, x + 0.3, yy, cw - 0.55, 0.7, [{"text": [
            {"text": mk, "size": 13, "bold": True, "color": (col if hi else MUTED)},
            {"text": it, "size": 11.8, "color": SLATE2, "line_spacing": 1.05}]}])
        yy += 0.74 if "\n" in it else 0.5
band = rrect(s, 0.7, 5.5, 11.96, 0.92, fill=SLATE, line=None, radius=0.13, shadow=True, grad=(H_SLATE, H_INDIGO_D), grad_ang=0)
tb(s, 1.05, 5.55, 11.3, 0.85, [
    {"text": [{"text": "Tư tưởng cốt lõi:  ", "size": 14.5, "bold": True, "color": RGBColor(0x7D,0xD3,0xFC)},
              {"text": "LLM chỉ là bộ xử lý ngôn ngữ có kiểm soát — KHÔNG phải nguồn tri thức pháp luật độc lập.", "size": 14.5, "color": WHITE}]}
], anchor=MSO_ANCHOR.MIDDLE)

# ============================================================================
# SLIDE 5 — CÔNG NGHỆ
# ============================================================================
s = slide()
header(s, "Công nghệ", "Công nghệ sử dụng", 5)
groups = [
    ("⚙️", "Backend & Điều phối", ["Python · FastAPI", "2 microservice", "LangGraph (agent)"]),
    ("🤖", "AI / LLM", ["Google Gemini 2.5", "Flash · Pro (verifier)", "Google Search · Tavily"]),
    ("🔍", "Truy hồi ngữ nghĩa", ["vietnamese-bi-encoder", "cross-encoder rerank", "sentence-transformers"]),
    ("🗄️", "Cơ sở dữ liệu", ["PostgreSQL (ứng dụng)", "MongoDB (văn bản)", "ChromaDB (vector)"]),
    ("📱", "Giao diện", ["Kotlin Multiplatform", "Compose (mobile)", "Next.js (admin)"]),
    ("🔑", "Hạ tầng khác", ["JWT Auth", "Cloudinary (lưu file)", "SSE / WebSocket"]),
]
cw, ch, gx, gy, x0, y0 = 3.82, 1.86, 0.24, 0.24, 0.7, 1.98
for i, (em, t, items) in enumerate(groups):
    r, c = divmod(i, 3)
    x = x0 + c * (cw + gx); y = y0 + r * (ch + gy)
    acc, tint = PAL[i % len(PAL)]
    achex = "%02X%02X%02X" % (acc[0], acc[1], acc[2])
    rrect(s, x, y, cw, ch, line=BORDER, radius=0.1, shadow=True, grad=("FFFFFF", tint), grad_ang=120)
    rrect(s, x, y, 0.12, ch, fill=acc, line=None, radius=0.0)
    icon(s, x + 0.28, y + 0.26, 0.52, em, achex, achex, size=17)
    tb(s, x + 0.95, y + 0.34, cw - 1.1, 0.5, [{"text": t, "size": 13.5, "bold": True, "color": acc}])
    tb(s, x + 0.95, y + 0.86, cw - 1.1, 0.95, [{"text": it, "size": 11, "color": SLATE2, "space_after": 3} for it in items])

# ============================================================================
# SLIDE 6 — KIẾN TRÚC TỔNG QUAN
# ============================================================================
s = slide()
header(s, "Kiến trúc", "Sơ đồ tổng quan hệ thống", 6)
box_node(s, 1.0, 2.2, 2.55, 0.92, "Mobile App", "Người dùng · KMP", accent=SKY, a_hex=H_SKY, tint="E8F6FE", icon_emoji="📱")
box_node(s, 1.0, 3.45, 2.55, 0.92, "Web Admin", "Quản trị · Next.js", accent=SKY, a_hex=H_SKY, tint="E8F6FE", icon_emoji="🖥️")
box_node(s, 4.8, 2.5, 3.05, 1.55, "Main Service · 8000", "Auth (JWT) · Chat ·\nTra cứu luật · Upload", accent=INDIGO, a_hex=H_INDIGO, tint="EEEDFE", icon_emoji="⚙️", line_w=2.0)
box_node(s, 9.35, 2.5, 3.05, 1.55, "RAG Service · 8001", "Agentic RAG (LangGraph) ·\nTruy hồi · Ingest", accent=VIOLET, a_hex=H_VIOLET, tint="F2EBFE", icon_emoji="🧠", line_w=2.0)
box_node(s, 4.8, 4.55, 1.46, 0.9, "PostgreSQL", "User · Chat", accent=GREEN, a_hex=H_GREEN, tint="E7F8F1", line_w=1.4)
box_node(s, 6.39, 4.55, 1.46, 0.9, "MongoDB", "Văn bản", accent=GREEN, a_hex=H_GREEN, tint="E7F8F1", line_w=1.4)
box_node(s, 9.35, 4.55, 1.46, 0.9, "ChromaDB", "Vector", accent=TEAL, a_hex=H_TEAL, tint="E5F6F3", line_w=1.4)
box_node(s, 10.94, 4.55, 1.46, 0.9, "Gemini · Web", "LLM · Tavily", accent=AMBER, a_hex=H_AMBER, tint="FDF3E2", line_w=1.4)
connect(s, 3.55, 2.66, 4.8, 3.0, color=SKY)
connect(s, 3.55, 3.91, 4.8, 3.55, color=SKY)
connect(s, 7.85, 3.27, 9.35, 3.27, color=INDIGO, w=2.4)
connect(s, 5.53, 4.05, 5.53, 4.55, color=GREEN)
connect(s, 7.12, 4.05, 7.12, 4.55, color=GREEN)
connect(s, 10.08, 4.05, 10.08, 4.55, color=TEAL)
connect(s, 11.67, 4.05, 11.67, 4.55, color=AMBER)
pill(s, 3.35, 2.18, "JWT", H_SKY, H_TEAL, h=0.3, size=9, pad=0.16)
pill(s, 8.0, 2.92, "X-API-Key", H_INDIGO, H_VIOLET, h=0.3, size=9, pad=0.16)
tb(s, 0.7, 6.1, 11.9, 0.5, [{"text": "Hai service Python (FastAPI) tách biệt — client gọi Main Service bằng JWT; Main Service gọi RAG Service nội bộ bằng API key.", "size": 11.5, "color": MUTED}])

# ============================================================================
# SLIDE 7 — TÍNH NĂNG CHÍNH
# ============================================================================
s = slide()
header(s, "Tính năng", "Các tính năng chính", 7)
feats = [
    ("💬", "Chat tư vấn (Agentic RAG)", "Hỏi đáp tự nhiên, trả lời có nguồn, hiển thị tiến trình realtime.", 0),
    ("🧭", "Tư vấn có hướng dẫn", "Hệ thống hỏi làm rõ tình huống → trả lời sát trường hợp người dùng.", 1),
    ("🔍", "Tra cứu & AI Search", "Duyệt danh mục luật, lọc theo lĩnh vực/năm, tìm kiếm ngữ nghĩa.", 2),
    ("📄", "Số hoá văn bản (OCR)", "Admin upload PDF → tự trích xuất điều luật → nạp kho tri thức.", 3),
    ("📊", "Quản trị & Dashboard", "Theo dõi tài liệu, tiến trình xử lý, thống kê tổng quan.", 4),
]
cw, ch, gx, gy, x0, y0 = 3.82, 1.96, 0.24, 0.24, 0.7, 1.98
for em, t, d, idx in feats:
    r, c = divmod(idx, 3)
    x = x0 + c * (cw + gx); y = y0 + r * (ch + gy)
    acc, tint = PAL[idx % len(PAL)]
    achex = "%02X%02X%02X" % (acc[0], acc[1], acc[2])
    rrect(s, x, y, cw, ch, line=BORDER, radius=0.1, shadow=True, grad=("FFFFFF", tint), grad_ang=120)
    icon(s, x + 0.3, y + 0.3, 0.66, em, achex, achex, size=22)
    tb(s, x + 1.12, y + 0.36, cw - 1.3, 0.7, [{"text": t, "size": 13.5, "bold": True, "color": acc, "line_spacing": 1.0}])
    tb(s, x + 0.3, y + 1.12, cw - 0.55, 0.8, [{"text": d, "size": 11.3, "color": SLATE2, "line_spacing": 1.14}])
# ô nhấn mạnh (vị trí thứ 6)
x = x0 + 2 * (cw + gx); y = y0 + 1 * (ch + gy)
rrect(s, x, y, cw, ch, line=None, radius=0.1, shadow=True, grad=(H_INDIGO, H_VIOLET), grad_ang=125)
icon(s, x + 0.3, y + 0.32, 0.6, "🛡️", "FFFFFF", "E8F6FE", size=20)
tb(s, x + 0.3, y + 1.02, cw - 0.6, 0.85, [
    {"text": "Điểm nhấn", "size": 12, "bold": True, "color": RGBColor(0x7D,0xD3,0xFC), "space_after": 4},
    {"text": "Mọi câu trả lời đều kèm nguồn pháp lý và đi qua bước kiểm chứng chống bịa đặt.", "size": 12, "color": WHITE, "line_spacing": 1.16}])

# ============================================================================
# SLIDE 8 — PIPELINE AGENTIC RAG
# ============================================================================
s = slide()
header(s, "Thuật toán trọng tâm · 1/2", "Pipeline xử lý câu hỏi (LangGraph)", 8)
nodes = [
    ("🛡️", "Guardrail", "Lọc câu lạc đề /\nđộc hại", INDIGO, H_INDIGO, "EEEDFE"),
    ("🔎", "Query Analysis", "Tách truy vấn tối ưu\n(nội bộ + web)", SKY, H_SKY, "E8F6FE"),
    ("🧠", "Agent", "Điều phối, ép gọi\nĐỦ 2 công cụ", VIOLET, H_VIOLET, "F2EBFE"),
    ("🛠️", "Tools", "retrieve_internal_law\nsearch_web_for_law", TEAL, H_TEAL, "E5F6F3"),
    ("✅", "Verifier", "Đối chiếu từng con số,\nxoá phần bịa", GREEN, H_GREEN, "E7F8F1"),
]
nw, nh, gap, y = 2.2, 1.5, 0.26, 2.3
x0 = 0.66
xs = []
for i, (em, t, sub, col, hx, tint) in enumerate(nodes):
    x = x0 + i * (nw + gap); xs.append(x)
    box_node(s, x, y, nw, nh, t, sub, accent=col, a_hex=hx, tint=tint, icon_emoji=em, line_w=2.0)
    if i < len(nodes) - 1:
        connect(s, x + nw, y + nh / 2, x + nw + gap, y + nh / 2, color=INDIGO, w=2.2)
connect(s, xs[3] + nw / 2, y, xs[2] + nw / 2, y, color=AMBER, w=1.8)
tb(s, xs[2] + nw - 0.3, y - 0.44, nw + gap, 0.3, [{"text": "lặp đến khi đủ nguồn (≤ 6 vòng)", "size": 9.5, "bold": True, "color": AMBER, "align": PP_ALIGN.CENTER}])
tb(s, xs[4], y + nh + 0.12, nw, 0.3, [{"text": "→ Trả câu trả lời", "size": 10.5, "bold": True, "color": GREEN, "align": PP_ALIGN.CENTER}])
band = rrect(s, 0.7, 4.65, 11.96, 1.68, line=BORDER, radius=0.1, shadow=True, grad=("FFFFFF", "F4F6FF"), grad_ang=120)
icon(s, 1.0, 4.86, 0.5, "⚡", H_INDIGO, H_VIOLET, size=15)
tb(s, 1.62, 4.88, 11.0, 0.4, [{"text": "Vì sao pipeline này quan trọng?", "size": 14, "bold": True, "color": INDIGO_D}])
tb(s, 1.05, 5.42, 11.5, 0.95, [
    {"text": [{"text": "• Đồ thị ÉP agent truy hồi đủ cả nguồn nội bộ lẫn web hiện hành trước khi trả lời — LLM không tự ý kết thúc sớm.\n", "size": 12.3, "color": SLATE2}], "space_after": 2},
    {"text": [{"text": "• ", "size": 12.3, "color": SLATE2},
              {"text": "Verifier (Gemini Pro)", "size": 12.3, "bold": True, "color": GREEN},
              {"text": " là lớp chống ảo giác: đối chiếu từng điều khoản & con số với kết quả tra cứu, chỉ giữ phần CÓ CĂN CỨ.", "size": 12.3, "color": SLATE2}]},
])

# ============================================================================
# SLIDE 9 — HYBRID RETRIEVAL
# ============================================================================
s = slide()
header(s, "Thuật toán trọng tâm · 2/2", "Truy hồi tri thức: Hybrid + chống luật hết hiệu lực", 9)
steps = [
    ("1", "Vector search", "Bi-encoder mã hoá truy vấn → tìm các đoạn gần nghĩa nhất trong ChromaDB.", SKY, H_SKY, "E8F6FE"),
    ("2", "Cross-encoder rerank", "Chấm lại độ liên quan và trộn điểm để xếp hạng chính xác hơn.", INDIGO, H_INDIGO, "EEEDFE"),
    ("3", "Year-boost", "Cộng điểm nhẹ cho văn bản mới hơn dựa trên năm ban hành.", TEAL, H_TEAL, "E5F6F3"),
    ("4", "Phát hiện xung đột thời gian", "Gắn nhãn ⛔ (cũ) / ✅ (mới) để LLM không dùng nhầm quy định hết hiệu lực.", AMBER, H_AMBER, "FDF3E2"),
]
y = 1.98
for i, (n, t, d, col, hx, tint) in enumerate(steps):
    h = 0.95
    yy = y + i * (h + 0.14)
    rrect(s, 0.7, yy, 7.7, h, line=BORDER, radius=0.1, shadow=True, grad=("FFFFFF", tint), grad_ang=120)
    nb = oval(s, 0.92, yy + 0.23, 0.5, fill=col, grad=(hx, hx), grad_ang=45)
    _center_text(nb, n, 17, WHITE)
    tb(s, 1.62, yy + 0.14, 6.7, 0.4, [{"text": t, "size": 13.5, "bold": True, "color": col}])
    tb(s, 1.62, yy + 0.5, 6.7, 0.45, [{"text": d, "size": 11.3, "color": SLATE2, "line_spacing": 1.06}])
rrect(s, 8.6, 1.98, 4.06, 2.0, line=None, radius=0.12, shadow=True, grad=(H_SLATE, H_INDIGO_D), grad_ang=125)
tb(s, 8.9, 2.18, 3.5, 0.4, [{"text": "Công thức điểm xếp hạng", "size": 12.5, "bold": True, "color": RGBColor(0x7D,0xD3,0xFC)}])
tb(s, 8.9, 2.66, 3.5, 1.2, [
    {"text": "score = 0.3 × vector", "size": 14, "bold": True, "color": WHITE, "space_after": 4},
    {"text": "      + 0.7 × cross-encoder", "size": 14, "bold": True, "color": WHITE, "space_after": 6},
    {"text": "      ± year-boost", "size": 12.5, "color": RGBColor(0xCB,0xD5,0xE1)},
])
rrect(s, 8.6, 4.12, 4.06, 2.28, line=INDIGO, line_w=1.4, radius=0.12, shadow=True, grad=("FFFFFF", "EEEDFE"), grad_ang=120)
icon(s, 8.88, 4.32, 0.46, "💡", H_AMBER, H_AMBER, size=14)
tb(s, 9.46, 4.36, 3.0, 0.4, [{"text": "Vì sao quan trọng?", "size": 12.5, "bold": True, "color": INDIGO_D}])
tb(s, 8.9, 4.92, 3.55, 1.4, [
    {"text": "Hai mô hình bù nhau: vector tìm nhanh ứng viên, cross-encoder chấm lại tinh hơn. Nhãn ⛔/✅ ngăn trích dẫn mức phạt từ văn bản đã bị thay thế.",
     "size": 11.5, "color": SLATE2, "line_spacing": 1.16}])

# ============================================================================
# SLIDE 10 — DEMO
# ============================================================================
s = slide()
header(s, "Demo", "Một số giao diện thực tế", 10)
mobiles = [("mobile_chat_thinking.png", "Chat · tiến trình xử lý"),
           ("mobile_chat_answer_1.png", "Chat · trả lời có nguồn"),
           ("mobile_guided_step2.png", "Tư vấn có hướng dẫn")]
mh = 4.2; x = 0.75
for fn, cap in mobiles:
    pic = pic_h(s, x, 1.98, mh, os.path.join(IMG_DIR, fn), cap)
    if pic: x += (pic.width / EMU_IN) + 0.18
rect(s, 7.28, 2.05, 0.02, 4.35, BORDER)
for fn, cap, ay in [("admin_dashboard.png", "Admin · bảng điều khiển", 1.98),
                    ("admin_upload_processing.png", "Admin · số hoá văn bản (OCR)", 4.38)]:
    pic_w(s, 7.65, ay, 4.45, os.path.join(IMG_DIR, fn), cap)

# ============================================================================
# SLIDE 11 — ĐÁNH GIÁ
# ============================================================================
s = slide()
header(s, "Đánh giá", "Kết quả đánh giá (bộ 60 câu hỏi)", 11)
metrics = [
    ("🎯", "~90%", "Độ chính xác câu hỏi factual", "Bộ N1 · 30 câu có đáp án chuẩn", GREEN, H_GREEN, "E7F8F1"),
    ("✅", "3.67/5", "Chất lượng câu hỏi mở", "Bộ N2 · 30 câu · 20/30 đạt ≥ 4đ", INDIGO, H_INDIGO, "EEEDFE"),
    ("📚", "~19", "Nguồn dẫn / câu trả lời", "Trích dẫn văn bản & nguồn web", SKY, H_SKY, "E8F6FE"),
]
cw, gap, x0, cy, ch = 3.82, 0.24, 0.7, 2.05, 2.55
for i, (em, big, t, sub, col, hx, tint) in enumerate(metrics):
    x = x0 + i * (cw + gap)
    rrect(s, x, cy, cw, ch, line=BORDER, radius=0.12, shadow=True, grad=("FFFFFF", tint), grad_ang=120)
    rrect(s, x, cy, cw, 0.14, fill=col, line=None, radius=0.0, grad=(hx, hx), grad_ang=0)
    icon(s, x + cw / 2 - 0.33, cy + 0.34, 0.66, em, hx, hx, size=22)
    tb(s, x + 0.2, cy + 1.12, cw - 0.4, 0.8, [{"text": big, "size": 40, "bold": True, "color": col, "align": PP_ALIGN.CENTER}])
    tb(s, x + 0.25, cy + 1.92, cw - 0.5, 0.4, [{"text": t, "size": 13, "bold": True, "color": SLATE, "align": PP_ALIGN.CENTER}])
    tb(s, x + 0.25, cy + 2.24, cw - 0.5, 0.4, [{"text": sub, "size": 10, "color": MUTED, "align": PP_ALIGN.CENTER}])
band = rrect(s, 0.7, 4.95, 11.96, 1.35, line=None, radius=0.12, shadow=True, grad=(H_TEAL, H_SKY), grad_ang=0)
tb(s, 1.05, 5.12, 11.3, 1.1, [
    {"text": [{"text": "Bộ test có nhóm câu “xung đột hiệu lực” ", "size": 13.5, "bold": True, "color": WHITE},
              {"text": "(luật cũ vs luật mới cùng chủ đề): hệ thống ưu tiên đúng quy định hiện hành.", "size": 13.5, "color": RGBColor(0xEC,0xFE,0xFF)}], "space_after": 5, "line_spacing": 1.12},
    {"text": "Đánh giá bán tự động · chấm theo rubric: độ chính xác · trích dẫn · xử lý xung đột hiệu lực.", "size": 11, "color": RGBColor(0xE0,0xF2,0xFE)},
], anchor=MSO_ANCHOR.MIDDLE)

# ============================================================================
# SLIDE 12 — HẠN CHẾ & HƯỚNG PHÁT TRIỂN
# ============================================================================
s = slide()
header(s, "Hạn chế & Hướng phát triển", "Hạn chế và hướng phát triển", 12)
rrect(s, 0.7, 1.98, 5.88, 4.4, line=BORDER, radius=0.1, shadow=True, grad=("FFFFFF", "FDF3E2"), grad_ang=120)
icon(s, 1.0, 2.22, 0.5, "⚠️", H_AMBER, H_AMBER, size=15)
tb(s, 1.62, 2.28, 4.8, 0.4, [{"text": "HẠN CHẾ", "size": 13.5, "bold": True, "color": AMBER}])
bullets(s, 1.0, 2.95, 5.4, 3.4, [
    "Độ trễ còn cao do gọi 2 công cụ + bước kiểm chứng (vài chục giây/câu).",
    "Phụ thuộc hạn mức (quota) API Gemini và nguồn tìm kiếm web.",
    "Phủ dữ liệu pháp luật chưa toàn diện mọi lĩnh vực.",
], size=13, gap=13, mk=AMBER)
rrect(s, 6.78, 1.98, 5.88, 4.4, line=None, radius=0.1, shadow=True, grad=(H_INDIGO, H_VIOLET), grad_ang=125)
icon(s, 7.08, 2.22, 0.5, "🚀", "FFFFFF", "E8F6FE", size=15)
tb(s, 7.7, 2.28, 4.8, 0.4, [{"text": "HƯỚNG PHÁT TRIỂN", "size": 13.5, "bold": True, "color": WHITE}])
bullets(s, 7.08, 2.95, 5.4, 3.4, [
    "Giảm độ trễ: stream từng token, chạy song song & rút bớt vòng lặp.",
    "Mở rộng kho văn bản và nhiều lĩnh vực pháp luật hơn.",
    "Fine-tune mô hình embedding tiếng Việt cho miền pháp lý.",
    "Đánh giá tự động quy mô lớn hơn, thêm kiểm thử hồi quy.",
], size=13, gap=12, color=WHITE, mk=RGBColor(0x7D, 0xD3, 0xFC))

# ============================================================================
# SLIDE 13 — KẾT LUẬN / CẢM ƠN
# ============================================================================
s = slide(decor_style="none")
bg = s.shapes[0]
panel = rrect(s, -2.0, 4.3, 8.0, 7.0, fill=INDIGO, line=None, radius=0.0, grad=("EEF1FE", "E9F4FE"), grad_ang=120)
DECOR.add(panel.shape_id)
o = oval(s, 10.6, -1.8, 4.2, grad=("EAECFE", "F3F0FE"), grad_ang=30, decor=True)
rect(s, 0, 0, SW, 0.16, INDIGO, grad=(H_INDIGO, H_SKY), grad_ang=0)
pill(s, 0.9, 0.7, "KẾT LUẬN", H_INDIGO, H_VIOLET, h=0.42, size=12)
tb(s, 0.9, 1.32, 11.4, 0.7, [{"text": "Tổng kết đóng góp của đề tài", "size": 28, "bold": True, "color": SLATE}])
rrect(s, 0.92, 2.04, 1.6, 0.08, fill=INDIGO, line=None, radius=0.5, grad=(H_INDIGO, H_SKY), grad_ang=0)
contribs = [
    ("⚖️", "Hệ thống Agentic RAG pháp luật tiếng Việt", "Trả lời có căn cứ, có trích dẫn, biết kiểm chứng tính cập nhật của văn bản.", INDIGO, H_INDIGO),
    ("🛡️", "Chống bịa đặt & chống luật hết hiệu lực", "Verifier đối chiếu từng con số; nhãn ⛔/✅ phân biệt quy định cũ – mới.", VIOLET, H_VIOLET),
    ("📱", "Sản phẩm hoàn chỉnh, đa nền tảng", "Microservice backend + ứng dụng di động + web quản trị.", SKY, H_SKY),
]
y = 2.4
for i, (em, t, d, col, hx) in enumerate(contribs):
    yy = y + i * 1.0
    icon(s, 0.95, yy, 0.66, em, hx, hx, size=21)
    tb(s, 1.85, yy - 0.02, 10.6, 0.5, [{"text": t, "size": 16, "bold": True, "color": SLATE}])
    tb(s, 1.85, yy + 0.44, 10.6, 0.5, [{"text": d, "size": 12, "color": SLATE2, "line_spacing": 1.08}])
band = rrect(s, 0.9, 5.65, 11.55, 1.1, line=None, radius=0.16, shadow=True, grad=(H_INDIGO, H_VIOLET), grad_ang=0)
tb(s, 1.0, 5.68, 11.35, 1.05, [
    {"text": "Cảm ơn Quý Thầy/Cô đã lắng nghe!", "size": 21, "bold": True, "color": WHITE, "align": PP_ALIGN.CENTER, "space_after": 2},
    {"text": f"Sinh viên thực hiện: {COVER['student']}", "size": 12.5, "color": RGBColor(0xE0, 0xE7, 0xFF), "align": PP_ALIGN.CENTER}],
   anchor=MSO_ANCHOR.MIDDLE)


# ============================================================================
# HIỆU ỨNG: transition + entrance (float-up/fade gợn sóng + wipe cho connector)
# ============================================================================
AUTO_PLAY = True     # True: nội dung tự chạy khi vào slide; False: bấm 1 lần mới chạy
DUR = 380            # thời lượng mỗi hiệu ứng (ms)
STEP = 70            # độ trễ so-le (ms)


def add_transition(slide_, kind="fade", spd="med"):
    sld = slide_._element
    tr = etree.SubElement(sld, qn("p:transition")); tr.set("spd", spd)
    etree.SubElement(tr, qn("p:" + kind))


def _cond(parent, delay, evt=None):
    if evt:
        c = etree.SubElement(parent, qn("p:cond")); c.set("delay", str(delay)); c.set("evt", evt)
        etree.SubElement(etree.SubElement(c, qn("p:tgtEl")), qn("p:sldTgt"))
    else:
        cl = etree.SubElement(parent, qn("p:stCondLst"))
        etree.SubElement(cl, qn("p:cond")).set("delay", str(delay))


def _spTgt(parent, spid):
    etree.SubElement(etree.SubElement(parent, qn("p:tgtEl")), qn("p:spTgt")).set("spid", str(spid))


def add_anim(slide_, items, dur=DUR, step=STEP):
    """items: list[(spid, effect)] với effect ∈ {'rise','wipe'}. Một build group, so-le."""
    if not items:
        return
    sld = slide_._element
    timing = etree.SubElement(sld, qn("p:timing"))
    tnLst = etree.SubElement(timing, qn("p:tnLst"))
    root = etree.SubElement(etree.SubElement(tnLst, qn("p:par")), qn("p:cTn"))
    root.set("id", "1"); root.set("dur", "indefinite"); root.set("restart", "never"); root.set("nodeType", "tmRoot")
    rch = etree.SubElement(root, qn("p:childTnLst"))
    seq = etree.SubElement(rch, qn("p:seq")); seq.set("concurrent", "1"); seq.set("nextAc", "seek")
    mseq = etree.SubElement(seq, qn("p:cTn")); mseq.set("id", "2"); mseq.set("dur", "indefinite"); mseq.set("nodeType", "mainSeq")
    mch = etree.SubElement(mseq, qn("p:childTnLst"))
    g = etree.SubElement(etree.SubElement(mch, qn("p:par")), qn("p:cTn")); g.set("id", "3"); g.set("fill", "hold")
    _cond(g, 0 if AUTO_PLAY else "indefinite")
    gch = etree.SubElement(g, qn("p:childTnLst"))
    nid = [4]

    def NID():
        v = nid[0]; nid[0] += 1; return str(v)

    for i, (spid, eff) in enumerate(items):
        c = etree.SubElement(etree.SubElement(gch, qn("p:par")), qn("p:cTn"))
        c.set("id", NID()); c.set("presetClass", "entr"); c.set("fill", "hold"); c.set("grpId", "0")
        if eff == "wipe":
            c.set("presetID", "22"); c.set("presetSubtype", "1")
        else:
            c.set("presetID", "10"); c.set("presetSubtype", "0")
        c.set("nodeType", "withEffect" if (AUTO_PLAY or i > 0) else "clickEffect")
        _cond(c, i * step)
        cch = etree.SubElement(c, qn("p:childTnLst"))
        # set visible
        st = etree.SubElement(cch, qn("p:set")); stb = etree.SubElement(st, qn("p:cBhvr"))
        stc = etree.SubElement(stb, qn("p:cTn")); stc.set("id", NID()); stc.set("dur", "1"); stc.set("fill", "hold")
        _cond(stc, 0)
        _spTgt(stb, spid)
        etree.SubElement(etree.SubElement(stb, qn("p:attrNameLst")), qn("p:attrName")).text = "style.visibility"
        etree.SubElement(etree.SubElement(st, qn("p:to")), qn("p:strVal")).set("val", "visible")
        if eff == "rise":
            # chuyển động trồi lên (ppt_y)
            an = etree.SubElement(cch, qn("p:anim")); an.set("calcmode", "lin"); an.set("valueType", "num")
            ab = etree.SubElement(an, qn("p:cBhvr")); ab.set("additive", "base")
            ac = etree.SubElement(ab, qn("p:cTn")); ac.set("id", NID()); ac.set("dur", str(dur))
            _spTgt(ab, spid)
            etree.SubElement(etree.SubElement(ab, qn("p:attrNameLst")), qn("p:attrName")).text = "ppt_y"
            tav = etree.SubElement(an, qn("p:tavLst"))
            for tm, val in (("0", "ppt_y+.04"), ("100000", "ppt_y")):
                t_ = etree.SubElement(tav, qn("p:tav")); t_.set("tm", tm)
                etree.SubElement(etree.SubElement(t_, qn("p:val")), qn("p:strVal")).set("val", val)
        # fade / wipe reveal
        ae = etree.SubElement(cch, qn("p:animEffect")); ae.set("transition", "in")
        ae.set("filter", "fade" if eff == "rise" else "wipe(right)")
        aeb = etree.SubElement(ae, qn("p:cBhvr"))
        aec = etree.SubElement(aeb, qn("p:cTn")); aec.set("id", NID()); aec.set("dur", str(dur))
        _spTgt(aeb, spid)
    # seq prev/next
    _cond(etree.SubElement(seq, qn("p:prevCondLst")), 0, evt="onPrev")
    _cond(etree.SubElement(seq, qn("p:nextCondLst")), 0, evt="onNext")
    # bldLst
    bld = etree.SubElement(timing, qn("p:bldLst"))
    for spid, _e in items:
        bp = etree.SubElement(bld, qn("p:bldP")); bp.set("spid", str(spid)); bp.set("grpId", "0")


def _items(slide_):
    out = []
    for sh in slide_.shapes:
        try:
            L = sh.left / EMU_IN; T = sh.top / EMU_IN; W = sh.width / EMU_IN; H = sh.height / EMU_IN
        except Exception:
            continue
        if sh.shape_id in DECOR:        # nền + trang trí
            continue
        if W >= SW * 0.95 and H >= SH * 0.95:
            continue
        if T >= 6.95:                   # footer
            continue
        eff = "wipe" if sh._element.tag.endswith("}cxnSp") else "rise"
        out.append((round(T, 3), round(L, 3), sh.shape_id, eff))
    out.sort()
    return [(x[2], x[3]) for x in out]


for idx, sld_ in enumerate(prs.slides, 1):
    add_transition(sld_, "fade", "med")
    add_anim(sld_, _items(sld_))

prs.save(OUT)
print("Saved:", OUT, "| slides:", len(prs.slides._sldIdLst))
